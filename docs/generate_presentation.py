"""
Generate beautified PowerPoint presentation for the Enterprise LOS/LMS System.

Version 3.0 - Complete system coverage with enhanced visual design.

Usage:
    pip install python-pptx
    python generate_presentation.py

Output: LOS_LMS_System_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ── Color Palette ──────────────────────────────────────────────────────────
NAVY = RGBColor(15, 32, 65)
DARK_BLUE = RGBColor(0, 51, 102)
BLUE = RGBColor(0, 102, 178)
LIGHT_BLUE = RGBColor(0, 150, 214)
TEAL = RGBColor(0, 166, 153)
ACCENT_ORANGE = RGBColor(230, 126, 34)
ACCENT_GREEN = RGBColor(39, 174, 96)
ACCENT_RED = RGBColor(192, 57, 43)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 247, 250)
MID_GRAY = RGBColor(100, 116, 139)
DARK_GRAY = RGBColor(51, 65, 85)
TEXT_DARK = RGBColor(30, 41, 59)
BODY_TEXT = RGBColor(55, 65, 81)


def _add_gradient_bg(slide, color_top=NAVY, color_bottom=DARK_BLUE):
    """Add a solid dark background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color_top


def _add_accent_bar(slide, y=0, height=Inches(0.06), color=LIGHT_BLUE):
    """Add a thin accent bar across the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), y, Inches(10), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_side_strip(slide, color=TEAL, width=Inches(0.12)):
    """Add a vertical accent strip on the left."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), width, Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_bottom_bar(slide, text, color=NAVY):
    """Add a footer bar at the bottom."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(10), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = RGBColor(180, 190, 210)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_icon_circle(slide, x, y, size, color, label):
    """Add a colored circle with a label inside."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_stat_card(slide, x, y, w, h, number, label, color):
    """Add a statistic card with large number and label."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.line.width = Pt(1)

    # Number
    num_box = slide.shapes.add_textbox(x, y + Inches(0.15), w, Inches(0.55))
    nf = num_box.text_frame
    nf.paragraphs[0].text = number
    nf.paragraphs[0].font.size = Pt(28)
    nf.paragraphs[0].font.bold = True
    nf.paragraphs[0].font.color.rgb = color
    nf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Label
    lbl_box = slide.shapes.add_textbox(x, y + Inches(0.65), w, Inches(0.35))
    lf = lbl_box.text_frame
    lf.paragraphs[0].text = label
    lf.paragraphs[0].font.size = Pt(10)
    lf.paragraphs[0].font.color.rgb = MID_GRAY
    lf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_title_slide(prs, title, subtitle):
    """Add a dark-themed title slide with accent elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_gradient_bg(slide, NAVY)

    # Accent bars
    _add_accent_bar(slide, y=Inches(2.2), height=Inches(0.04), color=TEAL)
    _add_accent_bar(slide, y=Inches(5.3), height=Inches(0.04), color=LIGHT_BLUE)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sb = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(1.2))
    sf = sb.text_frame
    sf.word_wrap = True
    p2 = sf.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(180, 200, 230)
    p2.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, number, title, subtitle=""):
    """Add a section divider slide with large number."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_gradient_bg(slide, DARK_BLUE)

    # Large number
    nb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(2), Inches(2.5))
    nf = nb.text_frame
    nf.paragraphs[0].text = f"{number:02d}"
    nf.paragraphs[0].font.size = Pt(96)
    nf.paragraphs[0].font.bold = True
    nf.paragraphs[0].font.color.rgb = TEAL

    # Title
    tb = slide.shapes.add_textbox(Inches(3), Inches(2.0), Inches(6.5), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(3), Inches(3.3), Inches(6.5), Inches(1))
        sf = sb.text_frame
        sf.word_wrap = True
        sf.paragraphs[0].text = subtitle
        sf.paragraphs[0].font.size = Pt(18)
        sf.paragraphs[0].font.color.rgb = RGBColor(160, 180, 210)

    # Line accent
    _add_accent_bar(slide, y=Inches(4.5), height=Inches(0.03), color=TEAL)

    _add_bottom_bar(slide, "Enterprise LOS/LMS Platform  |  Version 3.0  |  February 2026", NAVY)


def add_content_slide(prs, title, bullet_points, subtitle=None, footer=True):
    """Add a light content slide with accent strip and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_side_strip(slide, TEAL)

    # Title bar background
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Title text
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    y_offset = 1.2
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.4), Inches(1.15), Inches(9), Inches(0.4))
        sf = sb.text_frame
        sf.paragraphs[0].text = subtitle
        sf.paragraphs[0].font.size = Pt(14)
        sf.paragraphs[0].font.italic = True
        sf.paragraphs[0].font.color.rgb = MID_GRAY
        y_offset = 1.6

    # Bullet points
    cb = slide.shapes.add_textbox(Inches(0.6), Inches(y_offset), Inches(8.8), Inches(5))
    cf = cb.text_frame
    cf.word_wrap = True

    for i, point in enumerate(bullet_points):
        p = cf.paragraphs[0] if i == 0 else cf.add_paragraph()
        # Check if it's an indented line
        if point.startswith("   "):
            p.text = f"    {point.strip()}"
            p.font.size = Pt(14)
            p.font.color.rgb = MID_GRAY
            p.space_after = Pt(4)
        elif point == "":
            p.text = ""
            p.space_after = Pt(6)
        else:
            p.text = f"  {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(8)

    if footer:
        _add_bottom_bar(slide, "Enterprise LOS/LMS Platform  |  Version 3.0", NAVY)


def add_table_slide(prs, title, headers, rows, subtitle=None, footer=True):
    """Add a slide with styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_side_strip(slide, BLUE)

    # Title bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    y_offset = 1.3
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.4), Inches(1.15), Inches(9), Inches(0.35))
        sf = sb.text_frame
        sf.paragraphs[0].text = subtitle
        sf.paragraphs[0].font.size = Pt(13)
        sf.paragraphs[0].font.italic = True
        sf.paragraphs[0].font.color.rgb = MID_GRAY
        y_offset = 1.65

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    row_height = min(0.42, 4.5 / num_rows)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(y_offset),
        Inches(9.0), Inches(min(num_rows * row_height, 5.0))
    )
    table = table_shape.table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    if footer:
        _add_bottom_bar(slide, "Enterprise LOS/LMS Platform  |  Version 3.0", NAVY)


def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points, footer=True):
    """Add a two-column slide with card-style columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_side_strip(slide, TEAL)

    # Title bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    # Left card background
    lcard = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(4.3), Inches(5.2)
    )
    lcard.fill.solid()
    lcard.fill.fore_color.rgb = LIGHT_GRAY
    lcard.line.color.rgb = RGBColor(210, 218, 226)
    lcard.line.width = Pt(1)

    # Left title
    ltb = slide.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(4), Inches(0.45))
    ltf = ltb.text_frame
    ltf.paragraphs[0].text = left_title
    ltf.paragraphs[0].font.size = Pt(18)
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.color.rgb = DARK_BLUE

    # Left accent line
    la = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.95), Inches(1.5), Inches(0.03)
    )
    la.fill.solid()
    la.fill.fore_color.rgb = TEAL
    la.line.fill.background()

    # Left content
    lcb = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(3.9), Inches(4.0))
    lcf = lcb.text_frame
    lcf.word_wrap = True
    for i, pt in enumerate(left_points):
        p = lcf.paragraphs[0] if i == 0 else lcf.add_paragraph()
        if pt == "":
            p.text = ""
            p.space_after = Pt(4)
        elif pt.startswith("   "):
            p.text = f"    {pt.strip()}"
            p.font.size = Pt(12)
            p.font.color.rgb = MID_GRAY
            p.space_after = Pt(3)
        else:
            p.text = f"  {pt}"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(5)

    # Right card background
    rcard = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.3), Inches(4.3), Inches(5.2)
    )
    rcard.fill.solid()
    rcard.fill.fore_color.rgb = LIGHT_GRAY
    rcard.line.color.rgb = RGBColor(210, 218, 226)
    rcard.line.width = Pt(1)

    # Right title
    rtb = slide.shapes.add_textbox(Inches(5.3), Inches(1.45), Inches(4), Inches(0.45))
    rtf = rtb.text_frame
    rtf.paragraphs[0].text = right_title
    rtf.paragraphs[0].font.size = Pt(18)
    rtf.paragraphs[0].font.bold = True
    rtf.paragraphs[0].font.color.rgb = DARK_BLUE

    # Right accent line
    ra = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(1.95), Inches(1.5), Inches(0.03)
    )
    ra.fill.solid()
    ra.fill.fore_color.rgb = LIGHT_BLUE
    ra.line.fill.background()

    # Right content
    rcb = slide.shapes.add_textbox(Inches(5.3), Inches(2.1), Inches(3.9), Inches(4.0))
    rcf = rcb.text_frame
    rcf.word_wrap = True
    for i, pt in enumerate(right_points):
        p = rcf.paragraphs[0] if i == 0 else rcf.add_paragraph()
        if pt == "":
            p.text = ""
            p.space_after = Pt(4)
        elif pt.startswith("   "):
            p.text = f"    {pt.strip()}"
            p.font.size = Pt(12)
            p.font.color.rgb = MID_GRAY
            p.space_after = Pt(3)
        else:
            p.text = f"  {pt}"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(5)

    if footer:
        _add_bottom_bar(slide, "Enterprise LOS/LMS Platform  |  Version 3.0", NAVY)


def add_stats_slide(prs, title, stats, subtitle=None):
    """Add a slide with statistic cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_side_strip(slide, TEAL)

    # Title bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.4), Inches(1.15), Inches(9), Inches(0.35))
        sf = sb.text_frame
        sf.paragraphs[0].text = subtitle
        sf.paragraphs[0].font.size = Pt(13)
        sf.paragraphs[0].font.italic = True
        sf.paragraphs[0].font.color.rgb = MID_GRAY

    # Layout cards in rows
    y_start = Inches(1.7) if subtitle else Inches(1.4)
    cols_per_row = 4
    card_w = Inches(2.0)
    card_h = Inches(1.1)
    x_start = Inches(0.5)
    x_gap = Inches(0.2)
    y_gap = Inches(0.25)

    for i, (num, label, color) in enumerate(stats):
        row = i // cols_per_row
        col = i % cols_per_row
        x = x_start + col * (card_w + x_gap)
        y = y_start + row * (card_h + y_gap)
        _add_stat_card(slide, x, y, card_w, card_h, num, label, color)

    _add_bottom_bar(slide, "Enterprise LOS/LMS Platform  |  Version 3.0", NAVY)


def create_presentation():
    """Create the full beautified presentation with detailed module content."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ─────────────────────────────────────────────────
    add_title_slide(
        prs,
        "Enterprise Loan Origination &\nManagement System",
        "A Complete Lending Platform — Origination to Closure\n\n"
        "Retail  |  LAP  |  Co-Lending  |  Securitization  |  Investments\n"
        "FLDG  |  ECL Staging  |  Collections  |  Supply Chain\n\n"
        "Version 3.0  |  February 2026"
    )

    # ── Slide 2: Platform Stats ────────────────────────────────────────
    add_stats_slide(
        prs,
        "Platform at a Glance",
        [
            ("82", "ORM Models", DARK_BLUE),
            ("31", "Service Modules", TEAL),
            ("14", "API Routers", BLUE),
            ("453", "Automated Tests", ACCENT_GREEN),
            ("21,300+", "Lines of Code", DARK_BLUE),
            ("70+", "API Endpoints", TEAL),
            ("39", "Documentation Sections", BLUE),
            ("22", "Test Files", ACCENT_GREEN),
        ],
        subtitle="Complete enterprise lending lifecycle in a single platform"
    )

    # ── Slide 3: Technology Stack ──────────────────────────────────────
    add_table_slide(
        prs,
        "Technology Stack",
        ["Component", "Technology", "Version", "Purpose"],
        [
            ["API Framework", "FastAPI", ">= 0.110", "Async REST with auto OpenAPI docs"],
            ["ORM", "SQLAlchemy 2.0", ">= 2.0", "Type-safe Mapped[] declarative models"],
            ["Validation", "Pydantic v2", ">= 2.6", "Request/response schema validation"],
            ["Database", "SQLite / PostgreSQL", "3.x / 14+", "Dev: SQLite, Prod: PostgreSQL"],
            ["Configuration", "pydantic-settings", ">= 2.2", "Env-based config with .env files"],
            ["Migrations", "Alembic", ">= 1.13", "Schema versioning & migration"],
            ["Server", "Uvicorn (ASGI)", ">= 0.29", "Production ASGI server"],
            ["Testing", "pytest + httpx", ">= 8.0", "Unit, integration & API tests"],
            ["Docs Gen", "python-pptx", ">= 0.6", "Automated PPTX report generation"],
        ],
        subtitle="Modern, production-ready Python stack with zero JavaScript dependencies"
    )

    # ── Slide 4: Architecture ──────────────────────────────────────────
    add_content_slide(
        prs,
        "System Architecture — Layered Design",
        [
            "CLIENT LAYER:  Mobile  |  Web  |  External API  |  Swagger UI  |  Collection Software",
            "",
            "API LAYER — 14 FastAPI Routers, 70+ endpoints:",
            "   borrowers, loan_products, loan_applications, loan_accounts, loan_lifecycle",
            "   loan_partners, loan_participations, collaterals, collections",
            "   documents, holiday_calendars, benchmark_rates, health",
            "",
            "SERVICE LAYER — 31 stateless modules (Session passed explicitly):",
            "   Foundation (pure math, no DB): interest, frequency, calendar, floating_rate",
            "   Schedule Generation: schedule, advanced_schedule, fees",
            "   Lifecycle: payments, accrual, delinquency, restructure, prepayment, closure, lifecycle",
            "   Partners: co_lending, settlement, fldg, servicer_income",
            "   Institutional: securitization, investment, selldown",
            "   Risk & Compliance: ecl, par_report, eod",
            "   Platform: workflow, lap_workflow, rules_engine, kyc, supply_chain, collection, collateral",
            "",
            "DATA LAYER — 82 ORM models across 34 files  |  Numeric(18,2) for monetary fields",
            "   All monetary calculations use Decimal with ROUND_HALF_UP in service layer",
        ],
        subtitle="Each layer has clear boundaries — services are stateless and composable"
    )

    # ── Slide 5: Executive Summary ─────────────────────────────────────
    add_table_slide(
        prs,
        "System Capabilities — All Modules",
        ["Module", "Models", "Services", "Key Features"],
        [
            ["Loan Origination", "4", "3", "Borrower, KYC, Application, Products, 5-level Approval"],
            ["Loan Management", "4", "7", "EMI/Bullet/IO schedules, Payments, Accrual, DPD tracking"],
            ["LAP Collateral", "4", "2", "Property details, Valuations, LTV, Insurance, Legal verify"],
            ["Co-Lending", "5", "2", "80:20/90:10/100:0 splits, Partner ledger, Settlements"],
            ["FLDG", "3", "1", "First/Second Loss, Utilization, Recovery, Top-up threshold"],
            ["ECL (IFRS 9)", "6", "1", "Stage 1/2/3, PD x LGD x EAD, Scenario analysis, Provisions"],
            ["Collections", "5", "1", "Cases, Actions, PTP, Escalation rules, Dashboard"],
            ["Investments", "8", "1", "NCD, CP, Bond, G-Sec, YTM, MTM, Coupon schedules"],
            ["Securitization", "5", "1", "Pools, Tranches, Waterfall, Investor cash flows"],
            ["Selldown", "5", "1", "Full/Partial transfer, Gain/Loss, Collection split"],
            ["Servicer Income", "5", "1", "Fees, Excess spread, Withholding, TDS/GST"],
            ["Platform", "8", "5", "Workflow, Rules engine, KYC, Supply chain, EOD batch"],
        ],
        subtitle="82 models | 31 services | 70+ endpoints | 453 tests — all production-ready"
    )

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 1: LOAN ORIGINATION
    # ══════════════════════════════════════════════════════════════════════
    add_section_slide(prs, 1, "Loan Origination (LOS)",
                      "Borrower onboarding, KYC, credit assessment, product selection")

    # ── Borrower & KYC Detail ──────────────────────────────────────────
    add_two_column_slide(
        prs,
        "Borrower Management & KYC",
        "Borrower Model Fields",
        [
            "external_id — Unique business identifier",
            "first_name, last_name, date_of_birth",
            "email (unique), phone, gender",
            "pan_number (unique), aadhaar_number",
            "address_line1/2, city, state, pincode",
            "employment_type: salaried/self/business",
            "monthly_income, employer_name",
            "kyc_status: pending/verified/rejected",
            "risk_category: low/medium/high",
            "",
            "API: POST/GET/PATCH /borrowers",
            "Search by: email, PAN, status",
        ],
        "KYC & Credit Bureau",
        [
            "KYC Verification Types:",
            "   aadhaar, pan, voter_id, passport",
            "   driving_license, address_proof",
            "",
            "Verification Status Flow:",
            "   pending --> verified / rejected",
            "   Match score (0-100) tracking",
            "",
            "Credit Bureau Reports:",
            "   CIBIL, Experian, Equifax, CRIF",
            "   credit_score, report_date",
            "   risk_flags (JSON), is_latest flag",
            "",
            "KYC Requirements per product",
        ],
    )

    # ── Loan Products ──────────────────────────────────────────────────
    add_table_slide(
        prs,
        "Loan Products — 6 Configured Types",
        ["Product", "Code", "Rate Type", "Schedule", "Day Count", "Tenure", "Key Config"],
        [
            ["Home Loan", "HOME", "Floating", "EMI", "ACT/365", "12-360 mo", "base 8.5%, grace 15d"],
            ["Personal Loan", "PL", "Fixed", "EMI", "30/360", "6-60 mo", "base 14%, penalty 2%"],
            ["Business Loan", "BL", "Floating", "EMI", "ACT/365", "12-84 mo", "base 12%, grace 7d"],
            ["Vehicle Loan", "VL", "Fixed", "EMI", "30/360", "12-84 mo", "base 9.5%, fee 1%"],
            ["Gold Loan", "GL", "Fixed", "Bullet", "ACT/365", "3-12 mo", "base 11%, LTV 75%"],
            ["LAP", "LAP", "Floating", "EMI", "ACT/365", "12-240 mo", "base 10%, 5-lvl approve"],
        ],
        subtitle="Each product: code, name, currency, interest_rate_type, schedule_type, "
                 "min/max_tenure, processing_fee_rate, penalty_rate, grace_days"
    )

    # ── Loan Application ───────────────────────────────────────────────
    add_two_column_slide(
        prs,
        "Loan Application Lifecycle",
        "Application Model Fields",
        [
            "borrower_id (FK), product_id (FK)",
            "status: submitted/draft/data_entry/",
            "   under_review/approved/rejected/",
            "   referred_back/disbursed",
            "channel: branch/online/partner/DSA",
            "requested_amount, requested_tenure",
            "credit_score (from bureau)",
            "approved_amount, approved_rate",
            "approved_tenure, decision_reason",
            "decision_at (timestamp)",
            "branch_id, branch_name (LAP)",
            "submitted_at, created_at, updated_at",
        ],
        "Application Processing",
        [
            "Status Transitions:",
            "   submitted --> under_review",
            "   under_review --> approved/rejected",
            "   approved --> disbursed",
            "",
            "LAP Applications:",
            "   5-level workflow auto-syncs status",
            "   draft --> data_entry --> review",
            "   --> approved --> disbursed",
            "",
            "API Endpoints:",
            "   POST /loan-applications",
            "   GET (list with filters)",
            "   PATCH /{id}/status",
            "   4 workflow endpoints (LAP)",
        ],
    )

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 2: LOAN MANAGEMENT
    # ══════════════════════════════════════════════════════════════════════
    add_section_slide(prs, 2, "Loan Management (LMS)",
                      "Schedules, payments, accrual, delinquency, restructuring")

    # ── Loan Account Detail ────────────────────────────────────────────
    add_content_slide(
        prs,
        "Loan Account — Complete Model",
        [
            "Core: account_number (unique), application_id (FK), principal_amount, disbursement_date",
            "   interest_rate, interest_rate_type (fixed/floating), schedule_type (emi/bullet/interest_only)",
            "   tenure_months, start_date, maturity_date, day_count_convention, status",
            "",
            "Outstanding Balances (auto-refreshed on payment):",
            "   principal_outstanding, interest_outstanding, fees_outstanding, penalty_outstanding",
            "",
            "Floating Rate: benchmark_rate_id (FK), spread, floor_rate, cap_rate, rate_reset_frequency",
            "",
            "Delinquency: dpd (days past due), is_npa (bool), npa_date, npa_category",
            "   (substandard/doubtful_1/doubtful_2/doubtful_3/loss)",
            "",
            "ECL: ecl_stage (1/2/3), ecl_provision, sicr_flag",
            "",
            "Co-Lending: is_co_lent, co_lending_ratio, has_fldg_coverage",
            "",
            "Closure: closure_date, closure_type, settlement_amount",
            "Write-off: is_written_off, write_off_date, recovered_amount",
        ],
        subtitle="34 tracked fields — the central entity of the LMS"
    )

    # ── Financial Engine Detail ────────────────────────────────────────
    add_two_column_slide(
        prs,
        "Financial Calculation Engine",
        "Day-Count Conventions",
        [
            "30/360 — Corporate bonds (US/EU)",
            "   30-day months, 360-day year",
            "ACT/365 — Indian lending standard",
            "   Actual days, 365-day year",
            "ACT/360 — Money market instruments",
            "   Actual days, 360-day year",
            "ACT/ACT — Govt bonds (ISDA)",
            "   Actual days, actual year (365/366)",
            "",
            "Precision: Decimal(ROUND_HALF_UP)",
            "CENT = Decimal('0.01')",
            "RATE_PRECISION = Decimal('0.0000000001')",
            "_to_decimal() in every service",
        ],
        "Schedule Types & Formulas",
        [
            "EMI: P x r x (1+r)^n / ((1+r)^n - 1)",
            "   where r = annual_rate / 12 / 100",
            "Interest-Only: P x r per period",
            "   Principal bullet at maturity",
            "Bullet: All at maturity",
            "",
            "Advanced (via ScheduleConfiguration):",
            "Step-Up: EMI increases by step%",
            "Step-Down: EMI decreases by step%",
            "Balloon: large final (balloon_%)",
            "Moratorium: defer N months",
            "   capitalize / accrue / waive",
            "Custom JSON schedule",
        ],
    )

    # ── Payment Frequency & Calendar ───────────────────────────────────
    add_two_column_slide(
        prs,
        "Frequency, Calendar & Floating Rates",
        "Payment Frequencies",
        [
            "weekly: 52 periods/year",
            "biweekly: 26 periods/year",
            "monthly: 12 periods/year",
            "quarterly: 4 periods/year",
            "semi_annual: 2 periods/year",
            "annual: 1 period/year",
            "",
            "Holiday Calendar Model:",
            "   HolidayCalendar + Holiday entries",
            "   one_time or recurring holidays",
            "",
            "Business Day Adjustments:",
            "following, preceding",
            "modified_following, modified_preceding",
        ],
        "Floating Rate Engine",
        [
            "Benchmark Rates:",
            "   REPO (6.50%), MCLR (8.50%)",
            "   T-Bill (6.75%), SOFR, LIBOR",
            "   BenchmarkRateHistory tracks changes",
            "",
            "Effective Rate Calculation:",
            "   rate = benchmark + spread",
            "   rate = max(floor, min(cap, rate))",
            "",
            "Rate Reset:",
            "   Frequency: monthly/quarterly/annual",
            "   check_rate_reset_due(account)",
            "   apply_rate_reset() recalculates",
            "   get_rate_reset_schedule() plans",
        ],
    )

    # ── Payment Waterfall Detail ───────────────────────────────────────
    add_content_slide(
        prs,
        "Payment Waterfall — Detailed Allocation Logic",
        [
            "apply_payment(loan_account_id, amount, paid_at, channel, reference, db)",
            "",
            "Step 1: Load all unpaid installments ordered by due_date (oldest first)",
            "Step 2: For each installment, allocate in strict order:",
            "   a) Fees outstanding (penalties, late fees, processing fees)",
            "   b) Interest outstanding (accrued interest for the period)",
            "   c) Principal outstanding (reduces loan balance)",
            "Step 3: Create PaymentAllocation record per installment touched",
            "   Records: principal_allocated, interest_allocated, fees_allocated",
            "Step 4: Mark installment as 'paid' if fully settled, else 'partial'",
            "Step 5: Track unallocated_amount on Payment if excess remains",
            "",
            "_refresh_account_balances(account, db):",
            "   Recalculates principal/interest/fees_outstanding from all schedules",
            "",
            "compute_dpd(account, as_of_date):",
            "   DPD = (as_of_date - oldest_unpaid_due_date).days",
            "   Returns 0 if all installments paid, else days since earliest unpaid",
        ],
        subtitle="Waterfall: Fees --> Interest --> Principal | Multi-installment | Partial payment support"
    )

    # ── Accrual & Delinquency ──────────────────────────────────────────
    add_two_column_slide(
        prs,
        "Daily Accrual & Delinquency Engine",
        "Interest Accrual (accrual.py)",
        [
            "accrue_interest_daily(account, date):",
            "   daily = balance x rate / days_in_year",
            "   Creates InterestAccrual record",
            "   Tracks cumulative_accrued",
            "",
            "InterestAccrual Model:",
            "   accrual_date, opening_balance",
            "   interest_rate, benchmark_rate, spread",
            "   accrued_amount, cumulative_accrued",
            "   day_count_convention, status",
            "",
            "run_daily_accrual_batch(date):",
            "   Processes all active accounts",
            "   Skips already-accrued dates",
            "reset_cumulative_on_payment()",
        ],
        "Delinquency (delinquency.py)",
        [
            "DPD Buckets (RBI classification):",
            "   Standard: 0, SMA-0: 1-30",
            "   SMA-1: 31-60, SMA-2: 61-90",
            "   NPA Sub: 91-365, NPA Doubt: 366+",
            "",
            "Sticky NPA (RBI Rule):",
            "   Once NPA, stays NPA until DPD=0",
            "   Full cure required, not partial",
            "",
            "DelinquencySnapshot (daily):",
            "   dpd, bucket, overdue_principal",
            "   overdue_interest, overdue_fees",
            "   missed_installments count",
            "",
            "run_daily_delinquency_batch():",
            "   Updates DPD + NPA + snapshots",
        ],
    )

    # ── Lifecycle Operations ───────────────────────────────────────────
    add_two_column_slide(
        prs,
        "Loan Lifecycle Operations",
        "Restructuring (restructure.py)",
        [
            "Restructure Types:",
            "   rate_change — Lower interest rate",
            "   tenure_extension — More months",
            "   emi_rescheduling — New EMI amount",
            "   principal_haircut — Reduce balance",
            "",
            "LoanRestructure Model:",
            "   original_rate/tenure/emi",
            "   new_rate/tenure/emi",
            "   waived_interest, waived_penalty",
            "   approval_by, approval_date",
            "   status: pending/approved/applied",
            "",
            "Impact analysis via lifecycle.py",
        ],
        "Prepayment & Closure",
        [
            "Prepayment Actions:",
            "   reduce_emi — Lower EMI, same tenure",
            "   reduce_tenure — Same EMI, fewer months",
            "   foreclosure — Pay off entirely",
            "",
            "Penalty: prepayment_amount x penalty_rate",
            "",
            "Closure Types:",
            "   normal — All dues paid",
            "   settlement (OTS) — Agreed amount",
            "",
            "Write-Off Types:",
            "   full, partial, technical",
            "   WriteOff + WriteOffRecovery models",
            "   Recovery by: borrower, guarantor,",
            "   legal, collection_agency, asset_sale",
        ],
    )

    # ── Fee Management ─────────────────────────────────────────────────
    add_content_slide(
        prs,
        "Fee Management — Types, Charges & Waterfall",
        [
            "FeeType Model — Configurable fee definitions:",
            "   code, name, calculation_type (flat/percentage/slab), applies_to (loan/borrower/product)",
            "   charge_timing: upfront, on_disbursement, recurring, on_event, on_closure",
            "   is_taxable, tax_rate, waterfall_priority (allocation order in payment)",
            "",
            "ProductFee — Fee-to-Product linkage:",
            "   product_id (FK), fee_type_id (FK)",
            "   flat_amount or percentage_value, min_amount, max_amount",
            "   grace_days (before penalty applies), is_mandatory, is_waivable",
            "",
            "FeeCharge — Individual charges on accounts:",
            "   loan_account_id, fee_type_id, charge_date, due_date",
            "   amount, tax_amount, paid_amount, waived_amount",
            "   status: pending / paid / waived / partially_paid",
            "",
            "Fee Waterfall: Fees have highest priority in payment allocation",
            "   Allocated BEFORE interest and principal per installment",
        ],
        subtitle="3 models (FeeType, ProductFee, FeeCharge) | Flat, percentage & slab calculations"
    )

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 3: LAP COLLATERAL
    # ══════════════════════════════════════════════════════════════════════
    add_section_slide(prs, 3, "LAP Collateral Management",
                      "Property valuation, insurance, legal verification, LTV tracking")

    # ── Collateral Model Deep Dive ─────────────────────────────────────
    add_content_slide(
        prs,
        "Collateral Model — Complete Property Details",
        [
            "Property Classification:",
            "   property_type: residential/commercial/industrial/land/mixed_use",
            "   property_sub_type: flat/bungalow/row_house/office/warehouse/shop/plot/farm",
            "",
            "Address: address_line1/2, city, state, pincode, district, taluka, village",
            "Area: area_sqft, carpet_area_sqft, built_up_area_sqft, land_area_acres",
            "",
            "Ownership: owner_name, co_owner_name, ownership_type (sole/joint/leasehold/freehold)",
            "   title_deed_number, registration_date, registration_number",
            "   survey_number, cts_number, plot_number",
            "",
            "Valuation Snapshot (auto-updated): market_value, distress_value, realizable_value",
            "   ltv_ratio, valuation_date, valuer_name",
            "",
            "Legal: legal_status (pending/clear/issue_found), encumbrance_status, cersai_registration_id",
            "Charge: charge_type (first/second/pari_passu), charge_creation_date, charge_id",
            "Flags: status (active/released/under_review), is_primary_security, remarks",
        ],
        subtitle="Links to: loan_application_id (required), loan_account_id (optional, set after disbursement)"
    )

    # ── Valuation, Insurance, Legal ────────────────────────────────────
    add_table_slide(
        prs,
        "Collateral Sub-Models — Valuation, Insurance, Legal",
        ["Model", "Key Fields", "Auto-Update Behavior"],
        [
            ["CollateralValuation", "valuation_type (initial/periodic/re_valuation/distress),\n"
             "valuer_name, valuer_agency, market_value, realizable_value,\n"
             "distress_value, forced_sale_value, ltv_at_valuation, report_ref",
             "Updates parent: market_value,\nvaluation_date, ltv_ratio"],
            ["CollateralInsurance", "policy_number, provider, insured_value, premium_amount,\n"
             "start_date, expiry_date, renewal_date,\n"
             "status (active/expired/cancelled), is_assigned_to_lender",
             "Updates parent:\ninsurance snapshot fields"],
            ["CollateralLegalVerification", "verification_type: title_search, encumbrance_check,\n"
             "cersai_search, revenue_record, mutation_check\n"
             "status (pending/clear/issue_found/waived), findings",
             "If ALL clear --> parent\nlegal_status = 'clear'"],
        ],
        subtitle="13 API endpoints: CRUD + valuations + insurance + legal-verifications + LTV + link-account + summary"
    )

    # ── LTV & 5-Level Workflow ─────────────────────────────────────────
    add_two_column_slide(
        prs,
        "LTV Calculation & 5-Level Approval Workflow",
        "LTV = Outstanding / Market Value",
        [
            "calculate_ltv(collateral_id, db):",
            "   If account exists (post-disburse):",
            "     LTV = principal_outstanding / mkt_val",
            "   Else (pre-disburse):",
            "     LTV = approved_amount / mkt_val",
            "",
            "Auto-recalculates when:",
            "   New valuation added",
            "   Payment reduces outstanding",
            "   Collateral linked to account",
            "",
            "Collateral Summary API:",
            "   Full property + all valuations",
            "   + insurance + legal status",
            "   + current LTV ratio",
        ],
        "5-Level Approval (lap_workflow.py)",
        [
            "Stages (10 total):",
            "   draft --> branch_data_entry",
            "   --> branch_manager_review",
            "   --> regional_credit_review",
            "   --> central_credit_review",
            "   --> sanctioning_authority",
            "   --> approved --> disbursement",
            "   + referred_back, rejected",
            "",
            "Review stages can:",
            "   Approve --> next stage",
            "   Refer back --> branch_data_entry",
            "   Reject --> terminal (rejected)",
            "",
            "Auto-syncs LoanApplication.status",
            "Sets decision_at on approved/rejected",
        ],
    )

    # ── Document / Media ───────────────────────────────────────────────
    add_content_slide(
        prs,
        "Document & Media Management",
        [
            "Document Model — Enhanced with 11 additional fields for LAP:",
            "   collateral_id (FK) — Links media to specific collateral",
            "   media_type: document / photo / video",
            "   section: collateral_exterior, collateral_interior, site_visit, due_diligence,",
            "            title_deed, valuation_report, legal_opinion, kyc_documents",
            "",
            "File Metadata:",
            "   file_size_bytes, mime_type (image/jpeg, video/mp4, application/pdf)",
            "   thumbnail_path (for photo/video previews)",
            "",
            "Geo-Tagging (for site visit photos):",
            "   capture_latitude, capture_longitude (Numeric 10,7 — ~1cm precision)",
            "   captured_at (DateTime), captured_by (String — field agent name)",
            "   description (up to 1000 chars)",
            "",
            "API Endpoints: POST/GET /documents + GET /documents/by-collateral/{id}",
            "   GET /documents/by-section/{section} with optional collateral_id/media_type filters",
        ],
        subtitle="Backward-compatible: all new fields nullable/defaulted — existing docs unaffected"
    )

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 4: CO-LENDING & PARTNERSHIPS
    # ══════════════════════════════════════════════════════════════════════
    add_section_slide(prs, 4, "Co-Lending & Partnerships",
                      "Partner management, FLDG, settlements, servicer income")

    # ── Co-Lending Detail ──────────────────────────────────────────────
    add_content_slide(
        prs,
        "Co-Lending — Partner & Participation Models",
        [
            "LoanPartner Model:",
            "   partner_type: bank/nbfc/hfc/mfi/corporate",
            "   registration_number, rbi_license_number",
            "   default_share_percent (e.g., 80%), default_yield_rate",
            "   provides_fldg (bool), is_servicer (bool)",
            "   total_exposure_limit, current_exposure",
            "   bank_account, ifsc_code, tax_details (JSON)",
            "",
            "LoanParticipation Model (per-loan level):",
            "   loan_account_id, partner_id, participation_type (co_lending/assignment/securitization)",
            "   share_percent, interest_rate (partner's rate vs borrower rate)",
            "   principal_disbursed, principal_outstanding, total_collections",
            "   excess_spread_earned, fldg_arrangement_id, servicer_arrangement_id",
            "   Write-off tracking: written_off_principal/interest/fees, recovery_amount",
            "   ECL per partner: ecl_stage, ecl_provision",
            "",
            "PartnerLedger: entry_type (disbursement/collection/settlement/adjustment/fee/interest)",
            "PartnerSettlement + PartnerSettlementDetail: batch settlement processing",
        ],
        subtitle="3 partners in seed data: SBI (80% share), Bajaj Finance (90%), HDFC Bank (70%)"
    )

    # ── FLDG Deep Dive ─────────────────────────────────────────────────
    add_two_column_slide(
        prs,
        "FLDG — First Loss Default Guarantee",
        "FLDGArrangement Model",
        [
            "arrangement_code, name",
            "partner_id (FK), anchor_entity_id",
            "guarantee_type: first_loss / second_loss",
            "guarantee_form: cash_deposit / bank_guar",
            "",
            "Limits & Balances:",
            "   guarantee_percent, max_guarantee_amount",
            "   current_fldg_balance",
            "   total_utilized, total_recovered",
            "",
            "Triggers:",
            "   trigger_dpd (e.g., 90 days)",
            "   trigger_on_write_off (bool)",
            "   trigger_on_npa (bool)",
            "",
            "Top-up: requires_top_up (bool)",
            "   top_up_threshold_percent (e.g., 80%)",
        ],
        "Utilization & Recovery",
        [
            "FLDGUtilization Model:",
            "   arrangement_id, loan_account_id",
            "   utilization_type: default/write_off/npa",
            "   claimed_amount, approved_amount",
            "   claim_date, approval_date",
            "   status: claimed/approved/settled/rejected",
            "",
            "FLDGRecovery Model:",
            "   utilization_id, recovery_amount",
            "   recovery_date, recovery_source",
            "   Returned to FLDG pool balance",
            "",
            "FLDG Claim Calculation:",
            "   principal x share% + interest + fees",
            "   Capped at remaining FLDG balance",
            "",
            "Seed: SBI-Bajaj FLDG, Rs 2.25 Cr balance",
        ],
    )

    # ── Settlement & Servicer Income ───────────────────────────────────
    add_two_column_slide(
        prs,
        "Partner Settlement & Servicer Income",
        "Settlement Processing",
        [
            "PartnerSettlement:",
            "   settlement_code, partner_id",
            "   period_start_date, period_end_date",
            "   total_principal, total_interest",
            "   total_fees, servicer_fee",
            "   net_settlement_amount",
            "   status: draft/submitted/approved/paid",
            "",
            "PartnerSettlementDetail:",
            "   Per-loan breakdown of settlement",
            "   loan_account_id, principal/interest",
            "",
            "PartnerLedgerEntry:",
            "   Double-entry ledger per transaction",
            "   debit_amount, credit_amount, balance",
        ],
        "Servicer Income (servicer_income.py)",
        [
            "ServicerArrangement:",
            "   servicer_fee_percent (e.g., 0.50%)",
            "   fee_calculation_basis: outstanding/disbursed",
            "   excess_spread_share_percent",
            "",
            "ServicerIncomeAccrual:",
            "   Daily fee accrual per arrangement",
            "   accrual_amount, cumulative_accrued",
            "",
            "ExcessSpreadTracking:",
            "   borrower_rate - lender_rate = spread",
            "   spread x outstanding = income",
            "",
            "WithholdingTracker:",
            "   TDS rate, GST rate per payment",
            "   gross_amount - tds - gst = net",
            "   5 models, 25 tests",
        ],
    )

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 5: ECL & RISK
    # ══════════════════════════════════════════════════════════════════════
    add_section_slide(prs, 5, "ECL & Risk Management",
                      "IFRS 9 provisioning, delinquency classification, PAR reporting")

    # ── ECL Deep Dive ──────────────────────────────────────────────────
    add_content_slide(
        prs,
        "IFRS 9 ECL — Expected Credit Loss Engine",
        [
            "ECL = EAD x PD x LGD  (per loan, per stage)",
            "",
            "ECLConfiguration (per product/portfolio):",
            "   stage_1_pd (default 0.5%), stage_2_pd (5%), stage_3_pd (100%)",
            "   lgd (default 45%), cure_rate, discount_rate",
            "   scenario weights: base_weight (60%), adverse (30%), optimistic (10%)",
            "",
            "ECLStaging — Stage assignment logic (priority order):",
            "   1. Written-off account --> Stage 3",
            "   2. NPA (is_npa=True) --> Stage 3",
            "   3. DPD > 90 --> Stage 3",
            "   4. DPD > 30 --> Stage 2",
            "   5. Restructured loan --> Stage 2 (SICR)",
            "   6. SICR flag set --> Stage 2",
            "   7. Otherwise --> Stage 1",
            "",
            "ECLProvision: provision_date, ead, pd, lgd, ecl_amount, stage, scenario",
            "ECLMovement: tracks stage transitions (from_stage, to_stage, reason)",
            "ECLPortfolioSummary: aggregate by product/portfolio | 6 models, 29 tests",
        ],
        subtitle="Configurable PD/LGD per product | 3 scenarios (base/adverse/optimistic) with weights"
    )

    # ── Delinquency & NPA Detail ───────────────────────────────────────
    add_table_slide(
        prs,
        "Delinquency Classification & NPA Rules",
        ["Classification", "DPD Range", "ECL Stage", "NPA Sub-category", "Provisioning Action"],
        [
            ["Standard", "0", "Stage 1", "N/A", "12-month ECL (0.5% PD)"],
            ["SMA-0", "1 - 30", "Stage 1", "N/A", "12-month ECL, SMS alert"],
            ["SMA-1", "31 - 60", "Stage 2", "N/A", "Lifetime ECL (5% PD), calls"],
            ["SMA-2", "61 - 90", "Stage 2", "N/A", "Lifetime ECL, field visit"],
            ["NPA", "91 - 365", "Stage 3", "Substandard", "100% PD, legal notice"],
            ["NPA", "366 - 730", "Stage 3", "Doubtful-1", "100% PD, recovery action"],
            ["NPA", "731 - 1095", "Stage 3", "Doubtful-2", "100% PD, asset seizure"],
            ["NPA", "1096+", "Stage 3", "Loss", "100% PD, write-off candidate"],
        ],
        subtitle="RBI Sticky NPA: Once NPA, stays NPA until DPD returns to 0 (full cure). "
                 "Partial payments do NOT cure NPA."
    )

    # ── PAR Report ─────────────────────────────────────────────────────
    add_content_slide(
        prs,
        "PAR Report — Portfolio at Risk Analysis",
        [
            "par_report.py generates comprehensive portfolio analytics with 60+ variables:",
            "",
            "LOS Variables (per application):",
            "   Borrower demographics, income, employment, credit score, KYC status",
            "   Product type, requested/approved amounts, channel, decision reason",
            "   Application-to-approval conversion, time-to-decision metrics",
            "",
            "LMS Variables (per account):",
            "   Outstanding balances (principal, interest, fees), DPD, bucket classification",
            "   NPA status & category, ECL stage & provision amount, SICR flag",
            "   Collection case status, collateral coverage & LTV ratio",
            "   Co-lending partner, FLDG coverage, restructure history",
            "   Payment behavior: last payment date, total paid, missed installments",
            "",
            "API Endpoints (4 on /loan-applications):",
            "   GET /loan-applications/{id}/par-variables — Single application LOS+LMS",
            "   GET /loan-applications/par-report — Portfolio-wide PAR with filters",
            "   GET /loan-applications/{id}/par-los — LOS variables only",
            "   GET /loan-applications/{id}/par-lms — LMS variables only",
        ],
        subtitle="Comprehensive risk analytics for regulatory reporting and portfolio monitoring"
    )

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 6: INSTITUTIONAL FEATURES
    # ══════════════════════════════════════════════════════════════════════
    add_section_slide(prs, 6, "Institutional Features",
                      "Fixed income investments, securitization, selldown")

    # ── Investment Detail ──────────────────────────────────────────────
    add_content_slide(
        prs,
        "Investment Module — Fixed Income Instruments",
        [
            "4 Instrument Types:",
            "   NCD — Non-Convertible Debentures (quarterly coupon, 3-5yr tenure)",
            "   CP — Commercial Paper (zero coupon, discounted, 3-12 month)",
            "   Bond — Corporate/Government (semi-annual coupon, 5-10yr)",
            "   G-Sec — Government Securities (semi-annual coupon, 10-30yr, sovereign risk)",
            "",
            "8 Models: InvestmentIssuer, InvestmentProduct, Investment,",
            "   InvestmentCouponSchedule, InvestmentAccrual, InvestmentValuation,",
            "   InvestmentTransaction, InvestmentPortfolioSummary",
            "",
            "Key Calculations:",
            "   Coupon Schedule: auto-generated based on coupon_frequency and tenure",
            "   YTM (Yield to Maturity): Newton-Raphson iterative solver",
            "      Solves: Price = Sum(CF_i / (1+y)^t_i) for y",
            "   MTM (Mark-to-Market): current_price vs book_value --> unrealized P&L",
            "   Daily Accrual: face_value x coupon_rate / days_in_year",
            "",
            "Seed: 4 investments — Tata Motors NCD, Reliance CP, HDFC Bond, GoI G-Sec 2034",
        ],
        subtitle="InvestmentIssuer: 4 issuers with credit ratings (AAA, A1+, AA+, Sovereign)"
    )

    # ── Securitization Detail ──────────────────────────────────────────
    add_two_column_slide(
        prs,
        "Securitization — Pools, Tranches & Cash Flows",
        "Pool Management",
        [
            "SecuritizationPool:",
            "   pool_code, pool_name",
            "   pool_type: ABS / MBS / CLO",
            "   total_pool_value, weighted_avg_rate",
            "   weighted_avg_tenure, cut_off_date",
            "   status: draft/active/closed/settled",
            "",
            "PoolLoan: loan_account_id --> pool_id",
            "   original_balance, current_balance",
            "   inclusion_date, exclusion_date",
            "",
            "PoolInvestment: investment_id --> pool",
            "   For investment-backed pools",
        ],
        "Tranche Waterfall & Investors",
        [
            "Investor Model:",
            "   investor_code, name, investor_type",
            "   pool_id, tranche (senior/mezzanine/equity)",
            "   investment_amount, yield_rate",
            "   principal_outstanding, status",
            "",
            "Waterfall Priority:",
            "   1. Senior tranche (lowest risk/yield)",
            "   2. Mezzanine tranche",
            "   3. Equity tranche (first loss, highest yield)",
            "",
            "InvestorCashFlow:",
            "   cash_flow_date, cf_type",
            "   principal/interest/fee amounts",
            "   status: scheduled/paid/overdue",
        ],
    )

    # ── Selldown Detail ────────────────────────────────────────────────
    add_content_slide(
        prs,
        "Selldown — Loan & Investment Transfers",
        [
            "SelldownTransaction Model:",
            "   transaction_code, transaction_type (full_transfer/partial_transfer)",
            "   seller_entity, buyer_id (FK to SelldownBuyer)",
            "   asset_type: loan / investment, asset_id (loan_account_id or investment_id)",
            "   book_value, sale_price, gain_loss (auto: sale_price - book_value)",
            "   sale_yield (YTM at sale price), settlement_date",
            "   status: initiated/approved/settled/cancelled",
            "",
            "SelldownBuyer: buyer_code, name, buyer_type, credit_rating, bank_details",
            "",
            "SelldownSettlement: settlement amounts, fees, tax, net_amount",
            "",
            "SelldownCollectionSplit (post-selldown):",
            "   Defines how future collections are split between seller and buyer",
            "   seller_share_percent, buyer_share_percent",
            "   Applies to: principal, interest, fees separately",
            "",
            "SelldownPortfolioSummary: aggregate metrics by asset type/buyer/period",
        ],
        subtitle="Seed: 3 buyers — Mutual Fund (AAA), Insurance Co (AA+), Pension Fund (AAA)"
    )

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 7: COLLECTIONS
    # ══════════════════════════════════════════════════════════════════════
    add_section_slide(prs, 7, "Collection Management",
                      "Case lifecycle, actions, PTP, escalation, dashboard")

    # ── Collection Case Detail ─────────────────────────────────────────
    add_content_slide(
        prs,
        "Collection Case Management — Full Lifecycle",
        [
            "CollectionCase Model:",
            "   case_number (auto: COL-NNNNNN), loan_account_id (FK)",
            "   status: open --> in_progress --> resolved / written_off --> closed",
            "   priority: low / medium / high / critical",
            "   assigned_to (collector name), assigned_queue (team queue)",
            "   dpd_at_open, overdue_amount_at_open (captured from loan account)",
            "   last_action_date, next_action_date, follow_up_count",
            "   resolution_type: paid / settlement / restructured / write_off",
            "   resolution_date, opened_date, closed_date",
            "",
            "open_collection_case(loan_account_id, assigned_to, assigned_queue, priority, db):",
            "   Auto-generates case_number from count, captures current DPD and overdue amounts",
            "",
            "update_case_status(case_id, status, resolution_type, db):",
            "   Validates transitions: open-->in_progress-->resolved-->closed, open-->written_off",
            "   Sets resolution_date on resolved/written_off, prevents invalid transitions",
        ],
        subtitle="Status validation: Cannot transition resolved-->in_progress or closed-->any"
    )

    # ── Collection Actions & PTP ───────────────────────────────────────
    add_two_column_slide(
        prs,
        "Collection Actions & Promise to Pay",
        "CollectionAction Model",
        [
            "case_id (FK), action_type:",
            "   call, sms, email, letter,",
            "   field_visit, legal_notice",
            "",
            "performed_by, action_date",
            "outcome: contacted / not_reachable /",
            "   promise_made / refused / partial_pay",
            "follow_up_required (bool)",
            "next_action_date",
            "notes (free text)",
            "",
            "Auto-behaviors on log_action():",
            "   Updates case.last_action_date",
            "   Updates case.next_action_date",
            "   If case 'open' --> auto 'in_progress'",
        ],
        "PromiseToPay (PTP) Model",
        [
            "case_id (FK)",
            "promise_date, payment_due_date",
            "promised_amount",
            "",
            "Status Lifecycle:",
            "   pending --> kept / broken / partial",
            "",
            "On update (update_promise_status):",
            "   actual_payment_date",
            "   actual_amount (may < promised)",
            "   notes (reason for broken/partial)",
            "",
            "Seed Data Examples:",
            "   SMA-1 case: 3 actions (call/sms/visit)",
            "   NPA case: 4 actions + legal notice",
            "   PTP: Rs 10,000 broken, Rs 5,000 partial",
        ],
    )

    # ── Escalation & Dashboard ─────────────────────────────────────────
    add_two_column_slide(
        prs,
        "Escalation Rules & Collection Dashboard",
        "EscalationRule Model",
        [
            "name, trigger conditions (OR logic):",
            "   trigger_dpd >= threshold",
            "   trigger_bucket = SMA-1/SMA-2/NPA",
            "   trigger_amount >= threshold",
            "",
            "action_type: assign_queue, send_sms,",
            "   send_email, legal_notice, escalate",
            "action_config (JSON parameters)",
            "applies_to_product_id (optional filter)",
            "priority (evaluation order), is_active",
            "",
            "Seed Rules (5):",
            "   SMA-0 (DPD 1): SMS reminder",
            "   SMA-1 (DPD 31): Phone queue",
            "   SMA-2 (DPD 61): Field visit queue",
            "   NPA (DPD 90): Legal team queue",
            "   High Amount (>50K): Legal notice",
        ],
        "Dashboard & Analytics",
        [
            "get_collection_dashboard(db):",
            "   total_cases, by status breakdown",
            "   cases by priority (low/med/high/crit)",
            "   resolution_rate = resolved / total",
            "",
            "get_overdue_accounts(min_dpd, max_dpd):",
            "   Query delinquent accounts for",
            "   automated case creation",
            "",
            "get_case_details(case_id):",
            "   Full case + all actions + all PTPs",
            "   + loan account summary",
            "   + collateral info (if LAP)",
            "",
            "14 API endpoints on /collections:",
            "Cases CRUD, Actions, PTP, Escalation,",
            "Dashboard, Overdue query",
        ],
    )

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 8: PLATFORM SERVICES
    # ══════════════════════════════════════════════════════════════════════
    add_section_slide(prs, 8, "Platform Services",
                      "Workflow engine, rules engine, KYC, supply chain, EOD")

    # ── Workflow & Rules Engine ─────────────────────────────────────────
    add_two_column_slide(
        prs,
        "Workflow Engine & Rules Engine",
        "Generic Workflow (workflow.py)",
        [
            "WorkflowDefinition:",
            "   code (unique), stages_json",
            "   transitions_json, initial_stage",
            "   final_stages_json, requirements",
            "",
            "WorkflowInstance:",
            "   definition_id, entity_type/id",
            "   current_stage, previous_stage",
            "   assigned_to, assigned_role",
            "   priority, sla_due_date, is_active",
            "",
            "WorkflowTransition: audit trail",
            "WorkflowTask: per-stage tasks",
            "",
            "Used by: LAP 5-level approval",
            "18 tests covering all transitions",
        ],
        "Rules Engine (rules_engine.py)",
        [
            "RuleSet:",
            "   code, name, entity_type",
            "   is_active, evaluation_mode",
            "   (all_rules / first_match)",
            "",
            "DecisionRule:",
            "   rule_set_id, rule_order",
            "   condition_json (flexible logic)",
            "   action_type: approve/reject/refer/",
            "     assign/notify/score_adjust",
            "   action_params (JSON)",
            "   stop_on_match (bool)",
            "",
            "RuleExecutionLog: audit of every",
            "   rule evaluation with input/output",
            "24 tests covering rule evaluation",
        ],
    )

    # ── KYC, Supply Chain, EOD ─────────────────────────────────────────
    add_two_column_slide(
        prs,
        "KYC, Supply Chain Finance & EOD Batch",
        "KYC & Supply Chain",
        [
            "KYC Module (kyc.py):",
            "   KYCVerification: aadhaar, pan, voter_id,",
            "     passport, driving_license, address",
            "   KYCRequirement: per product config",
            "   CreditBureauReport: CIBIL, Experian",
            "     credit_score, risk_flags, report_date",
            "",
            "Supply Chain Finance:",
            "   Counterparty: buyer/seller entities",
            "   CreditLimit: per counterparty limits",
            "   Invoice: invoice financing records",
            "   Bill discounting, factoring support",
            "   16 tests for supply chain logic",
        ],
        "EOD Batch (eod.py)",
        [
            "run_end_of_day(target_date, db):",
            "   Orchestrates daily batch processing",
            "",
            "Step 1: Daily Interest Accrual",
            "   accrue_interest_daily for all accounts",
            "",
            "Step 2: DPD Refresh",
            "   compute_dpd for all active accounts",
            "",
            "Step 3: Delinquency Classification",
            "   NPA evaluation with sticky rule",
            "   Delinquency snapshot creation",
            "",
            "Step 4: ECL Stage Evaluation",
            "   Re-stage based on updated DPD/NPA",
            "",
            "Returns: accounts_processed count",
            "   accrual/delinquency/ecl summaries",
        ],
    )

    # ══════════════════════════════════════════════════════════════════════
    # API & DATA MODEL REFERENCE
    # ══════════════════════════════════════════════════════════════════════
    add_section_slide(prs, 9, "API & Data Model Reference",
                      "Complete endpoint listing and model catalog")

    # ── API Endpoints Detail ───────────────────────────────────────────
    add_table_slide(
        prs,
        "API Endpoints — Complete Reference (70+)",
        ["Router", "Prefix", "Methods", "Endpoints"],
        [
            ["Health", "/health", "1", "GET /health"],
            ["Borrowers", "/borrowers", "4", "POST, GET /{id}, GET (list), PATCH /{id}"],
            ["Products", "/loan-products", "4", "POST, GET /{id}, GET (list), PATCH /{id}"],
            ["Applications", "/loan-applications", "12", "CRUD + 4 workflow + 4 PAR endpoints"],
            ["Accounts", "/loan-accounts", "4", "POST, GET /{id}, GET (list), PATCH /{id}"],
            ["Lifecycle", "/loan-lifecycle", "6", "restructure, prepayment, closure, write-off, schedule"],
            ["Collaterals", "/collaterals", "13", "CRUD + valuations + insurance + legal + LTV + summary"],
            ["Collections", "/collections", "14", "Cases, actions, PTP, escalation, dashboard, overdue"],
            ["Documents", "/documents", "5", "CRUD + by-collateral + by-section"],
            ["Partners", "/loan-partners", "4", "CRUD"],
            ["Participations", "/loan-participations", "3", "POST, GET /{id}, GET (list)"],
            ["Calendars", "/holiday-calendars", "4", "CRUD + holidays"],
            ["Rates", "/benchmark-rates", "4", "CRUD + rate history"],
        ],
        subtitle="All endpoints: Swagger UI at http://localhost:8000/docs | OpenAPI spec auto-generated"
    )

    # ── Data Model Catalog ─────────────────────────────────────────────
    add_table_slide(
        prs,
        "Data Model Catalog — 82 ORM Models",
        ["Category", "Count", "Models"],
        [
            ["Core Loan", "4", "Borrower, LoanProduct, LoanApplication, LoanAccount"],
            ["Schedule & Pay", "4", "RepaymentSchedule, Payment, PaymentAllocation, ScheduleConfig"],
            ["Fees", "3", "FeeType, ProductFee, FeeCharge"],
            ["Calendar/Rate", "4", "HolidayCalendar, Holiday, BenchmarkRate, BenchmarkRateHistory"],
            ["Accrual/Delinq", "2", "InterestAccrual, DelinquencySnapshot"],
            ["Lifecycle", "4", "LoanRestructure, Prepayment, WriteOff, WriteOffRecovery"],
            ["Collateral", "4", "Collateral, CollateralValuation, CollateralInsurance, CollateralLegalVerification"],
            ["Collection", "4", "CollectionCase, CollectionAction, PromiseToPay, EscalationRule"],
            ["Documents", "1", "Document (with 11 media/geo fields)"],
            ["Users/Auth", "2", "User, RolePermission"],
            ["Co-Lending", "5", "LoanPartner, LoanParticipation, LedgerEntry, Settlement, Detail"],
            ["FLDG", "3", "FLDGArrangement, FLDGUtilization, FLDGRecovery"],
            ["ECL", "6", "ECLConfig, ECLStaging, ECLProvision, ECLMovement, ECLPortfolio, ECLUpload"],
            ["Servicer", "5", "ServicerArrangement, Accrual, Distribution, ExcessSpread, Withholding"],
            ["Investment", "8", "Issuer, Product, Investment, Coupon, Accrual, Valuation, Txn, Portfolio"],
            ["Securitization", "5", "Pool, PoolLoan, PoolInvestment, Investor, InvestorCashFlow"],
            ["Selldown", "5", "Transaction, Buyer, Settlement, CollectionSplit, PortfolioSummary"],
        ],
        subtitle="SQLAlchemy 2.0 Mapped[] types | Numeric(18,2) for monetary | All with created_at/updated_at"
    )

    # ── Pydantic Schemas ───────────────────────────────────────────────
    add_table_slide(
        prs,
        "Pydantic v2 Schemas — 15 Schema Files",
        ["Schema File", "Classes", "Pattern"],
        [
            ["borrower.py", "BorrowerCreate, BorrowerRead, BorrowerUpdate", "Base + Create/Read/Update"],
            ["loan_product.py", "LoanProductCreate, LoanProductRead, LoanProductUpdate", "Base + CRUD variants"],
            ["loan_application.py", "LoanApplicationCreate/Read/Update + branch fields", "With workflow status"],
            ["loan_account.py", "LoanAccountCreate, LoanAccountRead, LoanAccountUpdate", "34-field Read schema"],
            ["repayment_schedule.py", "RepaymentScheduleRead", "Read-only (auto-generated)"],
            ["payment.py", "PaymentCreate, PaymentRead, PaymentAllocationSchema", "Create + Read + nested"],
            ["collateral.py", "Collateral + Valuation + Insurance + Legal schemas", "4 entity groups, CRUD each"],
            ["collection.py", "Case + Action + PTP + EscalationRule schemas", "4 entity groups"],
            ["document.py", "DocumentCreate, DocumentRead (with 11 media fields)", "Enhanced for LAP"],
            ["holiday_calendar.py", "HolidayCalendarCreate/Read, HolidaySchema", "Nested holidays"],
            ["benchmark_rate.py", "BenchmarkRateCreate/Read, RateHistorySchema", "With rate history"],
        ],
        subtitle="All Read schemas: ConfigDict(from_attributes=True) for ORM compatibility"
    )

    # ══════════════════════════════════════════════════════════════════════
    # TESTING, SEED DATA & DEPLOYMENT
    # ══════════════════════════════════════════════════════════════════════
    add_section_slide(prs, 10, "Testing, Seed Data & Deployment",
                      "453 tests, demo data, development & production setup")

    # ── Test Coverage Detail ───────────────────────────────────────────
    add_table_slide(
        prs,
        "Test Coverage — 453 Tests Across 22 Files",
        ["Test File", "Tests", "What It Covers"],
        [
            ["test_day_count.py", "41", "4 day-count conventions, year fractions, edge cases"],
            ["test_frequency.py", "44", "6 frequencies, due date generation, tenure calculations"],
            ["test_calendar.py", "38", "Business days, holidays, 4 adjustment types"],
            ["test_floating_rate.py", "21", "Benchmarks, effective rate, floor/cap, rate reset"],
            ["test_schedule.py", "22", "EMI, interest-only, bullet schedule generation"],
            ["test_advanced_schedule.py", "14", "Step-up, step-down, balloon, moratorium"],
            ["test_fees.py", "17", "Fee types, charges, waivers, waterfall priority"],
            ["test_accrual.py", "9", "Daily accrual, cumulative tracking, batch runs"],
            ["test_lifecycle.py", "26", "Restructure, prepayment, closure, write-off"],
            ["test_ecl.py", "29", "ECL staging, PD/LGD, scenarios, provisions, movements"],
            ["test_fldg.py", "21", "FLDG arrangement, utilization, recovery, claims"],
            ["test_servicer.py", "25", "Servicer fees, excess spread, TDS/GST, accrual"],
            ["test_collateral.py", "20", "Collateral CRUD, LTV, valuation, legal auto-status"],
            ["test_lap_workflow.py", "20", "10 stages, transitions, refer-back, rejection"],
            ["test_collection_api.py", "22", "Cases, actions, PTP, escalation, dashboard"],
            ["test_delinquency_par.py", "5", "Delinquency buckets, PAR report variables"],
        ],
        subtitle="Run: pytest (from project root) | pytest.ini sets pythonpath=backend"
    )

    # ── Demo Seed Data Detail ──────────────────────────────────────────
    add_table_slide(
        prs,
        "Demo Seed Data — Comprehensive Test Dataset",
        ["Entity", "Count", "Details"],
        [
            ["Borrowers", "10", "Rajesh, Priya, Amit, Sunita, Vikram + 5 LAP borrowers (Meera, Arjun...)"],
            ["Loan Products", "6", "Home (8.5%), Personal (14%), Business (12%), Vehicle (9.5%), Gold (11%), LAP (10%)"],
            ["Applications", "9", "5 general (various statuses) + 4 LAP (workflow stages)"],
            ["Loan Accounts", "6", "Current (0 DPD), SMA-0 (15 DPD), SMA-1 (45 DPD), SMA-2 (75 DPD), NPA (95 DPD)"],
            ["Collaterals", "4", "Mumbai 3BHK (Rs 1.2Cr), Gurgaon office (Rs 3.5Cr), Pune bungalow, Chennai shop"],
            ["Repayment Schedules", "6x", "12-36 installments per account with realistic amounts"],
            ["Payments", "Multiple", "On-time, partial, and missed payments creating DPD scenarios"],
            ["Collection Cases", "2", "SMA-1 case (3 actions, 1 PTP) + NPA case (4 actions, 2 PTPs)"],
            ["Escalation Rules", "5", "SMA-0 SMS, SMA-1 Phone, SMA-2 Visit, NPA Legal, High Amount"],
            ["Partners + FLDG", "3 + 1", "SBI (80%), Bajaj (90%), HDFC (70%) + FLDG Rs 2.25Cr balance"],
            ["Investments", "4", "Tata NCD (9%), Reliance CP (7.5%), HDFC Bond (8.5%), GoI GSec (7.25%)"],
            ["Benchmark Rates", "3", "REPO 6.50%, MCLR 8.50%, T-Bill 6.75% with history"],
        ],
        subtitle="Run: cd backend && python3 -m app.db.seed_data | Guard: skips if data exists"
    )

    # ── How to Run Detail ──────────────────────────────────────────────
    add_two_column_slide(
        prs,
        "How to Run — Development & Production",
        "Development Setup",
        [
            "# Clone and setup",
            "git clone <repo-url>",
            "cd backend",
            "python3 -m venv .venv",
            "source .venv/bin/activate",
            "pip install -r requirements.txt",
            "pip install -r requirements-dev.txt",
            "cp .env.example .env",
            "",
            "# Initialize & seed database",
            "python3 -m app.db.init_db",
            "python3 -m app.db.seed_data",
            "",
            "# Start server (auto-reload)",
            "uvicorn app.main:app --reload",
            "# --> http://localhost:8000/docs",
        ],
        "Production & Testing",
        [
            "# PostgreSQL (docker)",
            "docker compose up -d",
            "# Set DATABASE_URL in .env:",
            "   postgresql+psycopg2://los:los@",
            "   localhost:5432/los_lms",
            "",
            "# Production server",
            "uvicorn app.main:app \\",
            "   --host 0.0.0.0 --port 8000 \\",
            "   --workers 4",
            "",
            "# Run all 453 tests",
            "cd <project-root>  # NOT backend/",
            "python3 -m pytest",
            "python3 -m pytest tests/test_schedule.py",
            "python3 -m pytest -v --tb=short",
        ],
    )

    # ── Domain Rules Reference ─────────────────────────────────────────
    add_table_slide(
        prs,
        "Key Domain Rules & Formulas",
        ["Rule", "Formula / Definition", "Service"],
        [
            ["EMI", "P x r x (1+r)^n / ((1+r)^n - 1)", "schedule.py"],
            ["Payment Waterfall", "Fees --> Interest --> Principal (per installment)", "payments.py"],
            ["DPD", "(as_of_date - oldest_unpaid_due_date).days", "payments.py"],
            ["NPA", "DPD >= 90; Sticky until DPD = 0 (RBI rule)", "delinquency.py"],
            ["ECL", "EAD x PD x LGD; Stage 1=12mo, 2/3=lifetime", "ecl.py"],
            ["LTV", "loan_outstanding / market_value x 100", "collateral.py"],
            ["FLDG Claim", "principal x share% + interest + fees", "fldg.py"],
            ["YTM", "Newton-Raphson iterative: Price = Sum(CF/(1+y)^t)", "investment.py"],
            ["Selldown G/L", "Gain/Loss = sale_price - book_value", "selldown.py"],
            ["Effective Rate", "max(floor, min(cap, benchmark + spread))", "floating_rate.py"],
            ["Excess Spread", "(borrower_rate - lender_rate) x outstanding", "servicer_income.py"],
            ["Year Fraction", "ACT/365: actual_days / 365", "interest.py"],
        ],
        subtitle="All monetary: Decimal(ROUND_HALF_UP) | DB: Numeric(18,2) | Conversion: _to_decimal()"
    )

    # ── Roadmap ────────────────────────────────────────────────────────
    add_content_slide(
        prs,
        "Future Roadmap",
        [
            "Near-Term:",
            "   JWT/OAuth2 authentication & role-based access control (RBAC)",
            "   API rate limiting, throttling & request logging",
            "   Notification service (email, SMS, push) with templates",
            "   Webhook integrations for external collection/CRM systems",
            "",
            "Medium-Term:",
            "   React admin dashboard with real-time portfolio monitoring",
            "   Credit bureau API integration (CIBIL, Experian, CRIF)",
            "   E-signature & e-stamping integration for document execution",
            "   Document management with S3/MinIO cloud storage",
            "   Automated report generation (PDF/Excel) for regulators",
            "",
            "Long-Term:",
            "   ML-based credit scoring & early warning system",
            "   Fraud detection engine with behavioral analysis",
            "   Multi-tenant SaaS architecture with tenant isolation",
            "   Real-time analytics & BI dashboard with drill-down",
        ],
        subtitle="Planned enhancements for enterprise-grade deployment"
    )

    # ── Thank You ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_gradient_bg(slide, NAVY)
    _add_accent_bar(slide, y=Inches(2.8), height=Inches(0.04), color=TEAL)
    _add_accent_bar(slide, y=Inches(5.0), height=Inches(0.04), color=LIGHT_BLUE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(0.8))
    tf = tb.text_frame
    tf.paragraphs[0].text = "Thank You"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    info = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(7), Inches(1.8))
    inf = info.text_frame
    inf.word_wrap = True

    lines = [
        ("API Documentation:  http://localhost:8000/docs", RGBColor(180, 200, 230)),
        ("System Docs:  docs/LAP_LOS_LMS_Documentation.md", RGBColor(180, 200, 230)),
        ("", None),
        ("82 Models  |  31 Services  |  453 Tests  |  21,300+ Lines", RGBColor(100, 200, 200)),
        ("", None),
        ("Enterprise LOS/LMS v3.0 — Lending Made Simple", RGBColor(160, 170, 190)),
    ]

    for i, (line, color) in enumerate(lines):
        p = inf.paragraphs[0] if i == 0 else inf.add_paragraph()
        p.text = line
        p.font.size = Pt(16) if i < 2 else Pt(18) if i == 3 else Pt(15)
        if color:
            p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

    # Save
    output_path = "LOS_LMS_System_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    create_presentation()
